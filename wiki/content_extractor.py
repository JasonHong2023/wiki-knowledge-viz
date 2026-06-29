"""
Hermes Content Extractor — Multi-format content extraction for Wiki import.

Converts PDF, PPT, Excel, and image files to Markdown text for analysis.
支援 PDF、PPT、Excel、圖片等非 Markdown 格式轉換為可分析的 Markdown 文字。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Callable, Optional

_log = logging.getLogger(__name__)


class ContentExtractor:
    """多格式內容提取器 / Multi-format content extractor.
    
    Extracts text content from various file formats and converts to Markdown.
    從各種檔案格式提取文字內容並轉換為 Markdown。
    
    Supported formats / 支援格式:
        - .md (Markdown)
        - .pdf (PDF documents)
        - .pptx (PowerPoint presentations)
        - .xlsx (Excel spreadsheets)
        - .png, .jpg, .jpeg, .gif, .webp (Images with OCR)
    """
    
    # Format mapping / 格式映射
    SUPPORTED_FORMATS = {
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.txt': 'markdown',
        '.pdf': 'pdf',
        '.pptx': 'pptx',
        '.xlsx': 'xlsx',
        '.epub': 'epub',
        '.png': 'image',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.gif': 'image',
        '.webp': 'image',
    }
    
    def __init__(self, progress_callback: Optional[Callable[[str, int, str], None]] = None):
        """
        初始化提取器 / Initialize extractor.
        
        Args:
            progress_callback: 可選的進度回呼函式 / Optional progress callback function
                              Signature / 簽名：callback(step: str, progress: int, message: str)
        """
        self.progress_callback = progress_callback
    
    def extract(self, file_path: Path) -> str:
        """
        根據檔案格式提取內容 / Extract content based on file format.
        
        Args:
            file_path: 檔案路徑 / File path
            
        Returns:
            提取的 Markdown 文字 / Extracted Markdown text
            
        Raises:
            ValueError: 不支援的檔案格式 / Unsupported file format
            FileNotFoundError: 檔案不存在 / File not found
        """
        suffix = file_path.suffix.lower()
        format_type = self.SUPPORTED_FORMATS.get(suffix)
        
        if format_type is None:
            raise ValueError(f"Unsupported file format: {suffix}")
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Dispatch to appropriate extractor / 路由到對應的提取器
        if format_type == 'markdown':
            return self._extract_md(file_path)
        elif format_type == 'pdf':
            return self._extract_pdf(file_path)
        elif format_type == 'pptx':
            return self._extract_pptx(file_path)
        elif format_type == 'xlsx':
            return self._extract_xlsx(file_path)
        elif format_type == 'image':
            return self._extract_image(file_path)
        elif format_type == 'epub':
            return self._extract_epub(file_path)

        raise ValueError(f"Unsupported file format: {suffix}")
    
    def _extract_md(self, file_path: Path) -> str:
        """
        提取 Markdown 檔案內容 / Extract Markdown file content.
        
        Args:
            file_path: Markdown 檔案路徑 / Markdown file path
            
        Returns:
            原始 Markdown 內容 / Original Markdown content
        """
        _log.debug("Extracting Markdown: %s", file_path)
        return file_path.read_text(encoding='utf-8')
    
    def _extract_pdf(self, file_path: Path) -> str:
        """
        提取 PDF 檔案文字 / Extract text from PDF file.
        
        Uses PyMuPDF (fitz) to extract text from PDF documents.
        使用 PyMuPDF (fitz) 從 PDF 文件提取文字。
        
        Args:
            file_path: PDF 檔案路徑 / PDF file path
            
        Returns:
            提取的文字內容（Markdown 格式） / Extracted text content (Markdown format)
            
        Raises:
            ImportError: PyMuPDF 未安裝 / PyMuPDF not installed
        """
        _log.debug("Extracting PDF: %s", file_path)
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF (fitz) is required for PDF extraction. Install with: pip install PyMuPDF")
        
        doc = fitz.open(file_path)
        total_pages = len(doc)
        text_pages = []
        
        for page_num in range(total_pages):
            page = doc[page_num]
            text = page.get_text()
            text_pages.append(text)
            
            # 回報進度（如果有 callback 且頁數 > 20）
            if self.progress_callback and total_pages > 20:
                progress = int((page_num + 1) / total_pages * 100)
                self.progress_callback(
                    "content_extraction",
                    progress,
                    f"正在解析 PDF：第 {page_num + 1}/{total_pages} 頁"
                )
        
        doc.close()
        
        # 組合為 Markdown 格式（頁間用分隔線）
        return "\n\n---\n\n".join(text_pages)
    
    def _extract_pptx(self, file_path: Path) -> str:
        """
        提取 PowerPoint 簡報文字 / Extract text from PowerPoint presentation.
        
        Uses python-pptx to extract text from slides.
        使用 python-pptx 從投影片提取文字。
        
        Args:
            file_path: PPTX 檔案路徑 / PPTX file path
            
        Returns:
            提取的文字內容（Markdown 格式） / Extracted text content (Markdown format)
            
        Raises:
            ImportError: python-pptx 未安裝 / python-pptx not installed
        """
        _log.debug("Extracting PPTX: %s", file_path)
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX extraction. Install with: pip install python-pptx")
        
        prs = Presentation(file_path)
        total_slides = len(prs.slides)
        slide_texts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text_parts = []
            
            # 提取標題
            if slide.shapes.title and slide.shapes.title.text:
                slide_text_parts.append(f"# {slide.shapes.title.text}")
            
            # 提取內文（所有 shape 的文字）
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        slide_text_parts.append(shape.text)
            
            slide_text = "\n\n".join(slide_text_parts)
            slide_texts.append(f"## Slide {slide_num + 1}\n\n{slide_text}")
            
            # 回報進度（如果有 callback 且投影片數 > 30）
            if self.progress_callback and total_slides > 30:
                progress = int((slide_num + 1) / total_slides * 100)
                self.progress_callback(
                    "content_extraction",
                    progress,
                    f"正在解析 PPT：第 {slide_num + 1}/{total_slides} 張投影片"
                )
        
        # 組合為 Markdown 格式
        return "\n\n".join(slide_texts)
    
    def _extract_xlsx(self, file_path: Path) -> str:
        """
        提取 Excel 工作表資料 / Extract data from Excel spreadsheet.
        
        Uses openpyxl to extract data from worksheets.
        使用 openpyxl 從工作表提取資料。
        
        Args:
            file_path: XLSX 檔案路徑 / XLSX file path
            
        Returns:
            提取的資料（Markdown 表格格式） / Extracted data (Markdown table format)
            
        Raises:
            ImportError: openpyxl 未安裝 / openpyxl not installed
        """
        _log.debug("Extracting XLSX: %s", file_path)
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl is required for XLSX extraction. Install with: pip install openpyxl")
        
        wb = load_workbook(filename=file_path, read_only=True, data_only=True)
        total_sheets = len(wb.sheetnames)
        sheet_texts = []
        
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            sheet_text_parts = [f"## Sheet: {sheet_name}\n"]
            
            # 獲取最大行數和列數
            max_row = ws.max_row or 0
            max_col = ws.max_column or 0
            
            if max_row == 0 or max_col == 0:
                sheet_text_parts.append("_(空工作表)_")
            else:
                # 提取表頭（第一行）
                headers = []
                for col in range(1, max_col + 1):
                    cell_val = ws.cell(row=1, column=col).value
                    headers.append(str(cell_val) if cell_val is not None else "")
                
                # 建立 Markdown 表格
                if headers:
                    sheet_text_parts.append("| " + " | ".join(headers) + " |")
                    sheet_text_parts.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                # 提取資料行（限制為前 100 行）
                data_rows = min(max_row - 1, 100)
                for row_idx in range(2, min(max_row + 1, 102)):
                    row_data = []
                    for col in range(1, max_col + 1):
                        cell_val = ws.cell(row=row_idx, column=col).value
                        row_data.append(str(cell_val) if cell_val is not None else "")
                    sheet_text_parts.append("| " + " | ".join(row_data) + " |")
                
                # 如果行數超過 100，添加註記
                if max_row > 101:
                    sheet_text_parts.append(f"\n_(共 {max_row} 行，僅顯示前 100 行)_")
            
            sheet_text = "\n".join(sheet_text_parts)
            sheet_texts.append(sheet_text)
            
            # 回報進度（如果有 callback 且工作表數 > 5）
            if self.progress_callback and total_sheets > 5:
                progress = int((sheet_idx + 1) / total_sheets * 100)
                self.progress_callback(
                    "content_extraction",
                    progress,
                    f"正在解析 Excel：第 {sheet_idx + 1}/{total_sheets} 張工作表"
                )
        
        # 組合為 Markdown 格式
        return "\n\n".join(sheet_texts)
    
    def _extract_epub(self, file_path: Path) -> str:
        """Extract text from EPUB ebook file."""
        try:
            from ebooklib import epub, ITEM_DOCUMENT
        except ImportError:
            raise ImportError("ebooklib is required for EPUB extraction. Install with: pip install ebooklib")
        import html2text as _h2t

        book = epub.read_epub(str(file_path), options={"ignore_ncx": True})

        h = _h2t.HTML2Text()
        h.ignore_links = False
        h.ignore_images = True
        h.body_width = 0
        h.single_line_break = False

        parts: list[str] = []

        # Book metadata header
        meta_title = book.get_metadata("DC", "title")
        meta_creator = book.get_metadata("DC", "creator")
        if meta_title:
            parts.append(f"# {meta_title[0][0]}")
        if meta_creator:
            parts.append(f"**作者**: {meta_creator[0][0]}\n")

        # Collect spine item IDs (reading order)
        spine_ids = {item_id for item_id, _ in book.spine}

        for item in book.get_items_of_type(ITEM_DOCUMENT):
            if item.get_id() not in spine_ids:
                continue
            try:
                html_bytes = item.get_content()
                html_str = html_bytes.decode("utf-8", errors="replace")
                md = h.handle(html_str).strip()
                if md and len(md) > 80:  # skip nav/toc pages
                    parts.append(md)
            except Exception:
                continue

        return "\n\n---\n\n".join(parts)

    def _extract_image(self, file_path: Path) -> str:
        """
        提取圖片文字（OCR）與描述 / Extract text (OCR) and description from image.
        
        Uses pytesseract for OCR and LLM vision for image description.
        使用 pytesseract 進行 OCR，使用 LLM vision 進行圖片描述。
        
        Args:
            file_path: 圖片檔案路徑 / Image file path
            
        Returns:
            OCR 文字 + 圖片描述（Markdown 格式） / OCR text + image description (Markdown format)
            
        Raises:
            ImportError: pytesseract 或 Pillow 未安裝 / pytesseract or Pillow not installed
        """
        _log.debug("Extracting image: %s", file_path)
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError("pytesseract and Pillow are required for image extraction. Install with: pip install pytesseract pillow")
        
        # 開啟圖片
        img = Image.open(file_path)
        
        # 1. OCR 文字辨識（使用 Tesseract）
        ocr_text = ""
        try:
            # 使用繁體中文 + 英文
            ocr_text = pytesseract.image_to_string(img, lang='chi_tra+eng')
        except Exception as e:
            _log.warning("OCR failed for %s: %s", file_path, e)
            ocr_text = "_(OCR 辨識失敗，可能是不含文字的圖片)_\n"
        
        # 2. 圖片基本資訊
        img_info = f"**圖片資訊**:\n- 尺寸：{img.width} x {img.height}\n- 模式：{img.mode}\n- 格式：{file_path.suffix.upper()}\n"
        
        # 3. 如果有 LLM vision 支援，可以添加圖片描述（TODO：未來實作）
        vision_desc = ""
        # TODO: 整合 Hermes 的 auxiliary client 進行 LLM vision 分析
        # 目前先使用簡單的圖片類型判斷
        if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
            vision_desc = "\n**圖片類型**: 含透明通道的圖片\n"
        elif img.mode == 'L':
            vision_desc = "\n**圖片類型**: 灰階圖片\n"
        elif img.mode == '1':
            vision_desc = "\n**圖片類型**: 黑白二值圖片\n"
        else:
            vision_desc = "\n**圖片類型**: 彩色圖片\n"
        
        # 4. 組合為 Markdown 格式
        result_parts = [
            f"# 圖片內容：{file_path.name}",
            img_info,
            vision_desc,
            "\n## OCR 辨識文字\n",
            ocr_text.strip() if ocr_text.strip() else "_(未辨識到文字)_",
        ]
        
        return "\n\n".join(result_parts)